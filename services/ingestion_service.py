from __future__ import annotations
import logging
import shutil
from pathlib import Path
from uuid import UUID

from core.database import get_driver
from models.document import Document
from repositories.document_repo import DocumentRepository
from services.concept_engine import extract_and_init_document_concepts

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".md", ".txt", ".pdf", ".docx", ".pptx", ".doc", ".ppt"}

# 副檔名 → (reader_func, file_type_str)
_READERS = {
    ".pdf":  ("_read_pdf",  "pdf"),
    ".md":   ("_read_text", "md"),
    ".txt":  ("_read_text", "txt"),
    ".docx": ("_read_docx", "docx"),
    ".pptx": ("_read_pptx", "pptx"),
    ".doc":  ("_read_doc",  "doc"),
    ".ppt":  ("_read_ppt",  "ppt"),
}


async def ingest_file(file_path: str) -> Document:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"找不到檔案：{file_path}")

    suffix = path.suffix.lower()
    if suffix not in _READERS:
        raise ValueError(f"不支援的檔案格式：{suffix}（支援：{', '.join(SUPPORTED_EXTENSIONS)}）")

    reader_name, file_type = _READERS[suffix]
    content = globals()[reader_name](path)

    if not content.strip():
        raise ValueError(f"檔案內容為空或無法解析：{path.name}")

    title = path.stem
    repo = DocumentRepository(get_driver())
    doc = await repo.create(
        title=title, content=content, file_path=str(path), file_type=file_type
    )
    await extract_and_init_document_concepts(doc.id, content)
    return doc


async def ingest_directory(dir_path: str) -> tuple[list[Document], list[str]]:
    """
    批次匯入目錄下所有支援格式的文件。
    回傳 (成功列表, 失敗訊息列表)。
    """
    path = Path(dir_path)
    if not path.exists():
        raise FileNotFoundError(f"目錄不存在：{dir_path}")

    success: list[Document] = []
    errors: list[str] = []

    files = [f for f in path.rglob("*") if f.suffix.lower() in SUPPORTED_EXTENSIONS]
    logger.info(f"找到 {len(files)} 個支援格式的文件，開始匯入…")

    for f in files:
        try:
            doc = await ingest_file(str(f))
            success.append(doc)
            logger.info(f"✅ 匯入成功：{f.name}")
        except Exception as e:
            msg = f"❌ {f.name}：{e}"
            errors.append(msg)
            logger.warning(msg)

    logger.info(f"匯入完成：成功 {len(success)} 個，失敗 {len(errors)} 個")
    return success, errors


async def move_and_ingest(
    source_dir: str,
    target_dir: str,
    delete_on_success: bool = True,
) -> tuple[list[Document], list[str]]:
    """
    將 source_dir 的文件移動到 target_dir 後匯入知識庫。
    delete_on_success=True（預設）：移動成功後刪除來源檔案。
    """
    src = Path(source_dir)
    dst = Path(target_dir)
    dst.mkdir(parents=True, exist_ok=True)

    files = [f for f in src.glob("*") if f.suffix.lower() in SUPPORTED_EXTENSIONS and f.is_file()]
    logger.info(f"準備搬移 {len(files)} 個文件：{src} → {dst}")

    moved: list[Path] = []
    move_errors: list[str] = []

    for f in files:
        dest_path = dst / f.name
        # 若目標已存在相同名稱，加上流水號避免覆蓋
        counter = 1
        while dest_path.exists():
            dest_path = dst / f"{f.stem}_{counter}{f.suffix}"
            counter += 1
        try:
            shutil.move(str(f), str(dest_path))
            moved.append(dest_path)
            logger.info(f"已移動：{f.name} → {dest_path}")
        except Exception as e:
            move_errors.append(f"移動失敗 {f.name}：{e}")
            logger.warning(f"移動失敗 [{f.name}]: {e}")

    # 匯入已移動的檔案
    success: list[Document] = []
    ingest_errors: list[str] = list(move_errors)

    for dest_path in moved:
        try:
            doc = await ingest_file(str(dest_path))
            success.append(doc)
            logger.info(f"✅ 匯入成功：{dest_path.name}")
        except Exception as e:
            msg = f"❌ {dest_path.name}：{e}"
            ingest_errors.append(msg)
            logger.warning(msg)
            # 匯入失敗則搬回來源（不遺失）
            try:
                shutil.move(str(dest_path), str(src / dest_path.name))
                logger.info(f"已還原：{dest_path.name} 搬回 {src}")
            except Exception:
                pass

    return success, ingest_errors


# ── 各格式讀取函式 ────────────────────────────────────────────────────────────

def _read_text(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "big5", "gbk", "cp950"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def _read_docx(path: Path) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            parts.append(f"[第 {i} 頁]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _read_doc(path: Path) -> str:
    """舊版 .doc 透過 Windows COM（需安裝 Microsoft Office）。"""
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = False
        doc = word.Documents.Open(str(path.absolute()), ReadOnly=True)
        text = doc.Content.Text
        doc.Close(False)
        word.Quit()
        return text
    except Exception as e:
        raise RuntimeError(f".doc 讀取失敗（需安裝 Microsoft Office）：{e}") from e


def _read_ppt(path: Path) -> str:
    """舊版 .ppt 透過 Windows COM（需安裝 Microsoft Office）。"""
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True
        presentation = ppt_app.Presentations.Open(
            str(path.absolute()), ReadOnly=True, Untitled=False, WithWindow=False
        )
        parts: list[str] = []
        for i, slide in enumerate(presentation.Slides, 1):
            slide_texts: list[str] = []
            for shape in slide.Shapes:
                try:
                    if shape.HasTextFrame:
                        t = shape.TextFrame.TextRange.Text.strip()
                        if t:
                            slide_texts.append(t)
                except Exception:
                    pass
            if slide_texts:
                parts.append(f"[第 {i} 頁]\n" + "\n".join(slide_texts))
        presentation.Close()
        ppt_app.Quit()
        return "\n\n".join(parts)
    except Exception as e:
        raise RuntimeError(f".ppt 讀取失敗（需安裝 Microsoft Office）：{e}") from e
