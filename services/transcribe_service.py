from __future__ import annotations
import logging
import re
import time
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".md", ".txt", ".pdf", ".docx", ".pptx", ".doc", ".ppt",
    # 音訊 / 影片（Whisper 轉譯）
    ".mp3", ".mp4", ".wav", ".m4a", ".ogg", ".flac", ".webm", ".mkv", ".avi",
}

AUDIO_EXTENSIONS = {".mp3", ".mp4", ".wav", ".m4a", ".ogg", ".flac", ".webm", ".mkv", ".avi"}


def _default_staging_dir() -> Path:
    from core.config import settings
    p = Path(settings.workspace_dir) / "_staging"
    p.mkdir(parents=True, exist_ok=True)
    return p


# ── 結果型別 ──────────────────────────────────────────────────────────────────

@dataclass
class TranscribeResult:
    src_path: str
    txt_path: str | None = None
    success: bool = False
    error: str | None = None
    char_count: int = 0
    elapsed_seconds: float = 0.0


# ── 公開 API ──────────────────────────────────────────────────────────────────

def transcribe_file(
    src_path: str | Path,
    staging_dir: str | Path | None = None,
    *,
    overwrite: bool = False,
) -> TranscribeResult:
    """
    將單一原始檔案轉譯成 .txt 並存入 staging_dir。
    回傳 TranscribeResult；無論成功或失敗都不拋例外。
    """
    src = Path(src_path)
    result = TranscribeResult(src_path=str(src))
    t0 = time.time()

    if not src.exists():
        result.error = f"找不到檔案：{src}"
        return result

    suffix = src.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        result.error = f"不支援的格式：{suffix}"
        return result

    dst_dir = Path(staging_dir) if staging_dir else _default_staging_dir()
    dst_dir.mkdir(parents=True, exist_ok=True)
    txt_path = _unique_txt_path(dst_dir, src.stem, overwrite=overwrite)

    try:
        text = _extract_text(src)
    except Exception as e:
        result.error = str(e)
        result.elapsed_seconds = time.time() - t0
        return result

    if not text.strip():
        result.error = "提取內容為空，可能是圖片型或加密檔案"
        result.elapsed_seconds = time.time() - t0
        return result

    txt_path.write_text(text, encoding="utf-8")
    result.txt_path = str(txt_path)
    result.success = True
    result.char_count = len(text)
    result.elapsed_seconds = round(time.time() - t0, 2)
    logger.info(f"轉譯完成：{src.name} → {txt_path.name}（{len(text):,} 字，{result.elapsed_seconds}s）")
    return result


def transcribe_folder(
    folder_path: str | Path,
    staging_dir: str | Path | None = None,
    *,
    recursive: bool = False,
    overwrite: bool = False,
) -> list[TranscribeResult]:
    """
    批次轉譯資料夾內所有支援格式的原始檔案。
    recursive=True 時遞迴子目錄（預設不遞迴，避免誤入 _staging/ 或 _text/）。
    """
    folder = Path(folder_path)
    if not folder.exists():
        raise FileNotFoundError(f"資料夾不存在：{folder}")

    glob = folder.rglob("*") if recursive else folder.glob("*")
    files = [f for f in glob if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS]
    logger.info(f"批次轉譯：{folder}，找到 {len(files)} 個支援格式的檔案")

    results: list[TranscribeResult] = []
    for f in files:
        r = transcribe_file(f, staging_dir, overwrite=overwrite)
        results.append(r)

    ok = sum(1 for r in results if r.success)
    fail = len(results) - ok
    logger.info(f"批次轉譯完成：✅ {ok}，❌ {fail}")
    return results


# ── 內部工具 ──────────────────────────────────────────────────────────────────

def _unique_txt_path(dst_dir: Path, stem: str, *, overwrite: bool) -> Path:
    """確保目標路徑不衝突；overwrite=False 時加流水號。"""
    p = dst_dir / f"{stem}.txt"
    if overwrite or not p.exists():
        return p
    counter = 1
    while True:
        p = dst_dir / f"{stem}_{counter}.txt"
        if not p.exists():
            return p
        counter += 1


def _extract_text(path: Path) -> str:
    """依副檔名分派到對應的讀取函式。"""
    suffix = path.suffix.lower()
    match suffix:
        case ".txt" | ".md":
            return _read_text(path)
        case ".pdf":
            return _read_pdf(path)
        case ".docx":
            return _read_docx(path)
        case ".pptx":
            return _read_pptx(path)
        case ".doc":
            return _read_doc(path)
        case ".ppt":
            return _read_ppt(path)
        case s if s in AUDIO_EXTENSIONS:
            return _read_audio(path)
        case _:
            raise ValueError(f"不支援的格式：{suffix}")


# ── 文字清理 ──────────────────────────────────────────────────────────────────

def _sanitize(text: str) -> str:
    cleaned = "".join(
        ch for ch in text
        if ch in ("\n", "\r", "\t") or unicodedata.category(ch)[0] != "C"
    )
    lines = [line.rstrip() for line in cleaned.splitlines()]
    return re.sub(r"\n{3,}", "\n\n", "\n".join(lines)).strip()


# ── 各格式讀取 ────────────────────────────────────────────────────────────────

def _read_text(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "big5", "gbk", "cp950"):
        try:
            return _sanitize(path.read_text(encoding=enc))
        except (UnicodeDecodeError, LookupError):
            continue
    return _sanitize(path.read_text(encoding="utf-8", errors="replace"))


def _read_pdf(path: Path) -> str:
    text = ""

    # Layer 1：pypdf
    try:
        from pypdf import PdfReader
        pages = [p.extract_text() or "" for p in PdfReader(str(path)).pages]
        text = _sanitize("\n".join(pages))
    except Exception:
        pass

    # Layer 2：pdfminer 備援
    if len(text.strip()) < 50:
        try:
            from pdfminer.high_level import extract_text as _mine
            fallback = _sanitize(_mine(str(path)) or "")
            if len(fallback.strip()) > len(text.strip()):
                text = fallback
        except Exception as e:
            logger.debug(f"pdfminer 備援失敗 [{path.name}]: {e}")

    # Layer 3：OCR 備援
    if len(text.strip()) < 50:
        try:
            logger.info(f"啟動 OCR：{path.name}")
            ocr = _ocr_pdf(path)
            if len(ocr.strip()) > len(text.strip()):
                text = ocr
                logger.info(f"OCR 完成：{path.name}，{len(text):,} 字")
        except Exception as e:
            logger.warning(f"OCR 失敗 [{path.name}]: {e}")

    return text


_ocr_reader = None


def _get_ocr_reader():
    global _ocr_reader
    if _ocr_reader is None:
        from paddleocr import PaddleOCR
        logger.info("初始化 PaddleOCR（中文 + 英文）…")
        # lang='ch' 支援繁/簡中文與英文；show_log=False 抑制冗長輸出
        _ocr_reader = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False, use_gpu=True)
        logger.info("PaddleOCR 載入完成（GPU）")
    return _ocr_reader


def _ocr_pdf(path: Path) -> str:
    import fitz
    import numpy as np

    reader = _get_ocr_reader()
    doc = fitz.open(str(path))
    parts: list[str] = []
    for i, page in enumerate(doc, 1):
        # 2x 縮放提升 OCR 準確度
        mat = fitz.Matrix(2.0, 2.0)
        pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
        img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, 3)

        result = reader.ocr(img, cls=True)
        page_lines: list[str] = []
        if result and result[0]:
            for line in result[0]:
                if line and len(line) >= 2:
                    text, _conf = line[1]
                    if text.strip():
                        page_lines.append(text)

        if page_lines:
            parts.append(f"[第 {i} 頁]\n" + "\n".join(page_lines))
        logger.debug(f"OCR {i}/{len(doc)}: {path.name}")
    doc.close()
    return _sanitize("\n\n".join(parts))


def _read_docx(path: Path) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return _sanitize("\n".join(parts))


def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        slide_texts.append(row_text)
        if slide_texts:
            parts.append(f"[第 {i} 頁]\n" + "\n".join(slide_texts))

    text = _sanitize("\n\n".join(parts))

    # OCR fallback：圖片型投影片（純圖、無文字框）
    if len(text.strip()) < 50:
        try:
            logger.info(f"啟動 PPTX OCR：{path.name}")
            ocr_text = _ocr_pptx(path)
            if len(ocr_text.strip()) > len(text.strip()):
                text = ocr_text
                logger.info(f"PPTX OCR 完成：{path.name}，{len(text):,} 字")
        except Exception as e:
            logger.warning(f"PPTX OCR 失敗 [{path.name}]: {e}")

    return text


def _ocr_pptx(path: Path) -> str:
    """從 PPTX 圖片形狀提取文字（適用純圖投影片）。"""
    import io
    import numpy as np
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    reader = _get_ocr_reader()
    prs = Presentation(str(path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        page_lines: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    img_arr = np.array(img)
                    result = reader.ocr(img_arr, cls=True)
                    if result and result[0]:
                        for line in result[0]:
                            if line and len(line) >= 2:
                                text, _conf = line[1]
                                if text.strip():
                                    page_lines.append(text)
                except Exception as e:
                    logger.debug(f"PPTX 形狀 OCR 失敗 [{path.name} p{i}]: {e}")

        if page_lines:
            parts.append(f"[第 {i} 頁]\n" + "\n".join(page_lines))
        logger.debug(f"PPTX OCR {i}/{len(prs.slides)}: {path.name}")

    return _sanitize("\n\n".join(parts))


def _read_doc(path: Path) -> str:
    import pythoncom
    import win32com.client
    pythoncom.CoInitialize()
    word = None
    doc = None
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = False
        doc = word.Documents.Open(str(path.absolute()), ReadOnly=True)
        text = doc.Content.Text
        return _sanitize(text)
    except Exception as e:
        raise RuntimeError(f".doc 讀取失敗（需安裝 Microsoft Office）：{e}") from e
    finally:
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word is not None:
            try:
                word.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def _read_ppt(path: Path) -> str:
    import pythoncom
    import win32com.client
    pythoncom.CoInitialize()
    ppt_app = None
    presentation = None
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True
        presentation = ppt_app.Presentations.Open(
            str(path.absolute()), ReadOnly=True, Untitled=False, WithWindow=False
        )
        parts: list[str] = []
        for i, slide in enumerate(presentation.Slides, 1):
            slide_texts: list[str] = []
            for shape in slide.Shapes:
                try:
                    if shape.HasTextFrame:
                        t = shape.TextFrame.TextRange.Text.strip()
                        if t:
                            slide_texts.append(t)
                except Exception:
                    pass
            if slide_texts:
                parts.append(f"[第 {i} 頁]\n" + "\n".join(slide_texts))
        return _sanitize("\n\n".join(parts))
    except Exception as e:
        raise RuntimeError(f".ppt 讀取失敗（需安裝 Microsoft Office）：{e}") from e
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if ppt_app is not None:
            try:
                ppt_app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


# ── Whisper 語音/影片轉譯 ──────────────────────────────────────────────────────

_whisper_model = None


def _get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        try:
            from faster_whisper import WhisperModel
        except ImportError as e:
            raise RuntimeError(
                "未安裝 faster-whisper。請執行：pip install faster-whisper"
            ) from e
        from core.config import settings
        model_size = settings.whisper_model_size
        logger.info(f"初始化 Whisper 模型（{model_size}）…")
        _whisper_model = WhisperModel(model_size, device="auto", compute_type="auto")
        logger.info("Whisper 模型載入完成")
    return _whisper_model


def _read_audio(path: Path) -> str:
    """使用 faster-whisper 將音訊或影片轉譯為文字。"""
    logger.info(f"啟動 Whisper 語音轉譯：{path.name}")
    model = _get_whisper_model()
    segments, info = model.transcribe(str(path), beam_size=5)
    lang = getattr(info, "language", "unknown")
    logger.info(f"Whisper 偵測語言：{lang}（{path.name}）")
    parts = [seg.text.strip() for seg in segments if seg.text.strip()]
    return _sanitize(" ".join(parts))
